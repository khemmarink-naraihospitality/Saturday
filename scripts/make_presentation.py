from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette ────────────────────────────────────────────────────────────────────
BRAND   = RGBColor(0x6B, 0x4C, 0xC3)   # Purple brand
ACCENT  = RGBColor(0x00, 0xC8, 0x75)   # Green
WARN    = RGBColor(0xFD, 0xAB, 0x3D)   # Orange
DANGER  = RGBColor(0xE2, 0x44, 0x5C)   # Red
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1A, 0x17, 0x28)
LIGHT   = RGBColor(0xF4, 0xF2, 0xFF)
GRAY    = RGBColor(0x64, 0x74, 0x8B)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Helper functions ────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, word_wrap=True, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_pill(slide, text, left, top, width, height, bg_rgb, text_rgb=WHITE, font_size=13):
    rect = add_rect(slide, left, top, width, height, fill_rgb=bg_rgb)
    rect.adjustments[0] = 0.5   # rounded corners via adjustment (works in some renderers)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = text_rgb
    run.font.name = "Segoe UI"

def header_bar(slide, color=BRAND):
    add_rect(slide, 0, 0, W, Inches(0.1), fill_rgb=color)

def slide_title(slide, title, subtitle=None, title_color=DARK, sub_color=GRAY):
    add_text(slide, title,
             Inches(0.8), Inches(0.25), Inches(11.5), Inches(0.7),
             font_size=30, bold=True, color=title_color)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.45),
                 font_size=16, color=sub_color)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 1 — Cover
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill_rgb=DARK)
add_rect(sl, 0, 0, Inches(0.45), H, fill_rgb=BRAND)

# Logo placeholder circle
add_rect(sl, Inches(1.1), Inches(1.3), Inches(1.4), Inches(1.4), fill_rgb=BRAND)
add_text(sl, "NHG", Inches(1.1), Inches(1.45), Inches(1.4), Inches(0.8),
         font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "NHG Saturday.com",
         Inches(1.0), Inches(3.0), Inches(10), Inches(1.1),
         font_size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_text(sl, "ระบบบริหารจัดการโปรเจกต์ภายในองค์กร",
         Inches(1.0), Inches(4.05), Inches(10), Inches(0.6),
         font_size=22, color=LIGHT, align=PP_ALIGN.LEFT)

add_text(sl, "Narai Hospitality Group  •  Business Tech",
         Inches(1.0), Inches(6.5), Inches(10), Inches(0.5),
         font_size=13, color=GRAY, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 2 — ระบบนี้คืออะไร?
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill_rgb=WHITE)
header_bar(sl)
slide_title(sl, "ระบบนี้คืออะไร?", "NHG Saturday.com — Project Management Platform ขององค์กร")

bullets = [
    ("📋", "Board & Group",      "จัดกลุ่มงานเป็น Board และ Group ได้อย่างยืดหยุ่น"),
    ("✅", "Items & Sub-items",  "สร้าง Item งาน และ Sub-item ย่อยได้ไม่จำกัด"),
    ("👥", "Team Collaboration", "กำหนดสิทธิ์สมาชิก เชิญทีม ติดตามผู้รับผิดชอบ"),
    ("📊", "Multiple Views",     "ดูงานได้ 4 มุมมอง: Table, Timeline, Kanban, Calendar"),
]
for i, (icon, title, desc) in enumerate(bullets):
    y = Inches(1.65) + Inches(1.22) * i
    add_rect(sl, Inches(0.7), y, Inches(11.5), Inches(1.0),
             fill_rgb=LIGHT, line_rgb=BRAND, line_width=Pt(1))
    add_text(sl, icon,  Inches(0.9),  y + Inches(0.18), Inches(0.6), Inches(0.7), font_size=22)
    add_text(sl, title, Inches(1.55), y + Inches(0.05), Inches(3.0), Inches(0.45),
             font_size=15, bold=True, color=BRAND)
    add_text(sl, desc,  Inches(1.55), y + Inches(0.48), Inches(9.5), Inches(0.45),
             font_size=13, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 3 — 4 มุมมองหลัก
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill_rgb=WHITE)
header_bar(sl)
slide_title(sl, "4 มุมมองการทำงาน", "เปลี่ยนวิธีดูข้อมูลได้ทันทีโดยไม่กระทบข้อมูลเดิม")

views = [
    ("📋", "Main Table",  "ตารางหลัก — ดูงานทุกชิ้นพร้อมกัน\nกรอง / เรียง / Group by ได้ทันที",        BRAND),
    ("📅", "Timeline",   "แผนภูมิ Gantt — วางแผนระยะเวลา\nแบบ Day / Month / Year",                      ACCENT),
    ("🗂️", "Kanban",     "กระดานการ์ด — ลาก-วางเปลี่ยนสถานะ\nขยาย Sub-item ได้บนการ์ด",               WARN),
    ("📆", "Calendar",   "มุมมองปฏิทิน — เห็น Deadline\nของทุก Item ในเดือนนั้น",                       DANGER),
]
for i, (icon, name, desc, color) in enumerate(views):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + Inches(6.3) * col
    y = Inches(1.55) + Inches(2.6) * row
    add_rect(sl, x, y, Inches(5.9), Inches(2.35), fill_rgb=WHITE,
             line_rgb=color, line_width=Pt(2))
    add_rect(sl, x, y, Inches(5.9), Inches(0.55), fill_rgb=color)
    add_text(sl, f"{icon}  {name}", x + Inches(0.15), y + Inches(0.05),
             Inches(5.5), Inches(0.45), font_size=18, bold=True, color=WHITE)
    add_text(sl, desc, x + Inches(0.2), y + Inches(0.65),
             Inches(5.5), Inches(1.6), font_size=13, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 4 — ฟีเจอร์เด่น (Filter / Sort / Group by)
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill_rgb=WHITE)
header_bar(sl)
slide_title(sl, "ฟีเจอร์เด่น — ค้นหาและจัดกลุ่มข้อมูล",
            "เครื่องมือ Toolbar ที่ใช้ได้ทุก View")

features = [
    ("🔍 Search",    "ค้นหา Item ได้แบบ Real-time\nรองรับทุก Field",            BRAND),
    ("🔽 Filter",    "กรองตาม Group / Status / Person\nหรือ Dropdown ได้พร้อมกัน", ACCENT),
    ("↕️ Sort",      "เรียง Item ตาม Column ใดก็ได้\nA→Z หรือ Z→A",              WARN),
    ("⬛ Group by",  "จัดกลุ่มแบบ Dynamic ตาม Status\nหรือ Dropdown โดยอัตโนมัติ", DANGER),
    ("👁️ Hide",      "ซ่อน Item ที่ไม่ต้องการแสดง\nออกจาก View ชั่วคราว",        RGBColor(0xA2,0x5D,0xDC)),
    ("📤 Export",    "Export เป็น Excel หรือ CSV\nนำไปใช้ต่อได้ทันที",            GRAY),
]
for i, (title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.45) + Inches(4.2) * col
    y = Inches(1.55) + Inches(2.6) * row
    add_rect(sl, x, y, Inches(3.9), Inches(2.35), fill_rgb=LIGHT,
             line_rgb=color, line_width=Pt(1.5))
    add_text(sl, title, x + Inches(0.18), y + Inches(0.15),
             Inches(3.5), Inches(0.5), font_size=16, bold=True, color=color)
    add_text(sl, desc, x + Inches(0.18), y + Inches(0.65),
             Inches(3.55), Inches(1.5), font_size=12, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 5 — Kanban View เจาะลึก
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill_rgb=WHITE)
header_bar(sl, WARN)
slide_title(sl, "Kanban View", "บริหารงานแบบ Card — ลาก-วาง เปลี่ยนสถานะได้ทันที")

points = [
    "🎨  แถบสี Status บนการ์ด — รู้สถานะทันทีโดยไม่ต้องเปิด",
    "👤  แสดง Person / ผู้รับผิดชอบบนการ์ด พร้อม Avatar",
    "📎  ดู Sub-items ได้โดยกด Expand ไม่ต้องออกจาก Kanban",
    "🖱️  Drag & Drop การ์ดข้าม Column เพื่อเปลี่ยน Status / Group",
    "🔽  Filter by Person — กรองดูเฉพาะงานของคนนั้นได้ทันที",
    "⬛  พื้นหลัง Column สีเทา — แยกแยะ Card กับพื้นหลังชัดเจน",
]
for i, pt in enumerate(points):
    y = Inches(1.55) + Inches(0.82) * i
    add_rect(sl, Inches(0.7), y, Inches(11.6), Inches(0.7),
             fill_rgb=RGBColor(0xFF,0xF8,0xEE) if i % 2 == 0 else WHITE,
             line_rgb=WARN if i % 2 == 0 else RGBColor(0xE2,0xE8,0xF0),
             line_width=Pt(1))
    add_text(sl, pt, Inches(0.95), y + Inches(0.12),
             Inches(11.2), Inches(0.5), font_size=14, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 6 — Import Board
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill_rgb=WHITE)
header_bar(sl, ACCENT)
slide_title(sl, "Import Board จาก Excel", "นำเข้าข้อมูลจากไฟล์ .xlsx โดยตรง — แทนที่ข้อมูลเดิมทั้งหมด")

steps = [
    ("1", "เปิด Board ที่ต้องการ Import",   "เลือก Board ปลายทางก่อน แล้วเปิดเมนู ⋯ > Import Board"),
    ("2", "เลือกไฟล์ Excel",                 "รองรับ .xlsx, หลายไฟล์พร้อมกัน, หลาย Sheet ต่อไฟล์"),
    ("3", "Preview & เลือก Sheet",           "ระบบแสดง Groups / Items ที่พบ ให้เลือกก่อน Import"),
    ("4", "กด Import",                        "ระบบลบ Group เดิมทั้งหมด แล้ววาง Group/Item จาก Excel"),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(1.6) + Inches(1.3) * i
    add_rect(sl, Inches(0.65), y, Inches(0.7), Inches(0.7), fill_rgb=ACCENT)
    add_text(sl, num, Inches(0.65), y + Inches(0.07),
             Inches(0.7), Inches(0.55), font_size=22, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, Inches(1.5), y + Inches(0.04),
             Inches(5.0), Inches(0.45), font_size=15, bold=True, color=ACCENT)
    add_text(sl, desc,  Inches(1.5), y + Inches(0.44),
             Inches(10.5), Inches(0.45), font_size=13, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 7 — สิทธิ์และการจัดการสมาชิก
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill_rgb=WHITE)
header_bar(sl, RGBColor(0xA2,0x5D,0xDC))
slide_title(sl, "การจัดการสมาชิกและสิทธิ์", "ควบคุมว่าใครทำอะไรได้บ้างในแต่ละ Board")

roles = [
    ("Owner",   "เจ้าของ Board — สิทธิ์เต็ม รวมถึงลบ Board ได้",                  BRAND),
    ("Admin",   "บริหารสมาชิก เพิ่ม/ลบคน กำหนดสิทธิ์",                             ACCENT),
    ("Editor",  "แก้ไข Item, เพิ่ม Group, อัปเดตค่าต่างๆ",                         WARN),
    ("Member",  "เพิ่ม Item และ Update ได้ ไม่สามารถลบ Group",                       RGBColor(0xA2,0x5D,0xDC)),
    ("Viewer",  "ดูข้อมูลได้อย่างเดียว ไม่สามารถแก้ไขใดๆ",                          GRAY),
]
for i, (role, desc, color) in enumerate(roles):
    y = Inches(1.6) + Inches(1.0) * i
    add_rect(sl, Inches(0.7), y, Inches(2.2), Inches(0.75), fill_rgb=color)
    add_text(sl, role, Inches(0.7), y + Inches(0.1),
             Inches(2.2), Inches(0.55), font_size=16, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(3.0), y, Inches(9.5), Inches(0.75),
             fill_rgb=WHITE, line_rgb=RGBColor(0xE2,0xE8,0xF0), line_width=Pt(1))
    add_text(sl, desc, Inches(3.2), y + Inches(0.18),
             Inches(9.1), Inches(0.45), font_size=13, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 8 — Tips การใช้งาน
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill_rgb=WHITE)
header_bar(sl, WARN)
slide_title(sl, "Tips การใช้งานที่มีประสิทธิภาพ", "เทคนิคที่ช่วยให้ทำงานได้เร็วขึ้น")

tips = [
    ("💡", "ใช้ Group by Status",    "เปลี่ยนมุมมองตาม Status ได้ทันทีจาก Toolbar > Group by"),
    ("💡", "Filter by Person",       "กรองดูเฉพาะงานของตัวเองหรือเพื่อนร่วมทีม"),
    ("💡", "Sub-items",              "ใช้ Sub-item สำหรับงานย่อยแทนการสร้าง Board ใหม่"),
    ("💡", "Timeline View",          "ใช้ Month View สำหรับวางแผนรายเดือน, Day สำหรับงาน Sprint"),
    ("💡", "Export ก่อน Import",     "Export Board ปัจจุบันก่อนเสมอ เพื่อ Backup ข้อมูล"),
    ("💡", "Updates / Comments",     "บันทึกความคืบหน้าไว้ใน Updates ของแต่ละ Item"),
]
for i, (icon, title, desc) in enumerate(tips):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(6.3) * col
    y = Inches(1.55) + Inches(1.7) * row
    add_rect(sl, x, y, Inches(6.0), Inches(1.55), fill_rgb=RGBColor(0xFF,0xFB,0xEB),
             line_rgb=WARN, line_width=Pt(1))
    add_text(sl, icon + "  " + title, x + Inches(0.2), y + Inches(0.1),
             Inches(5.5), Inches(0.5), font_size=14, bold=True, color=WARN)
    add_text(sl, desc, x + Inches(0.2), y + Inches(0.6),
             Inches(5.6), Inches(0.8), font_size=12, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 9 — Thank You / Q&A
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill_rgb=DARK)
add_rect(sl, 0, 0, Inches(0.45), H, fill_rgb=BRAND)
add_rect(sl, Inches(0.45), Inches(3.3), W - Inches(0.45), Inches(0.06), fill_rgb=BRAND)

add_text(sl, "ขอบคุณที่รับฟัง 🙏",
         Inches(1.0), Inches(1.8), Inches(11), Inches(1.2),
         font_size=42, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_text(sl, "หากมีคำถามหรือต้องการสาธิตเพิ่มเติม",
         Inches(1.0), Inches(3.55), Inches(11), Inches(0.6),
         font_size=20, color=LIGHT, align=PP_ALIGN.LEFT)

add_text(sl, "ติดต่อ: Business Tech Team  |  NHG Saturday.com",
         Inches(1.0), Inches(6.3), Inches(11), Inches(0.5),
         font_size=14, color=GRAY, align=PP_ALIGN.LEFT)

# ── Save ───────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "NHG_Saturday_Presentation.pptx")
prs.save(out)
print(f"✅  Saved: {out}")
